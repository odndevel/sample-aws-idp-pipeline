import asyncio
from contextlib import contextmanager

import boto3
import requests
from botocore.config import Config
from nanoid import generate
from strands import Agent, tool
from strands.models import BedrockModel
from strands_tools import current_time, http_request
from strands_tools.code_interpreter import AgentCoreCodeInterpreter

from agents.constants import REPORT_MODEL_ID
from config import get_config

NANOID_ALPHABET = "0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ-_"


def get_image_tool(bucket_name: str, artifact_base_path: str, unsplash_key: str | None):
    """Create an image search and S3 upload tool.

    Args:
        bucket_name: S3 bucket name for uploading images
        artifact_base_path: Base path in S3 for storing images
        unsplash_key: Unsplash API access key (optional)

    Returns:
        Tool function for searching and downloading images
    """

    @tool
    def search_and_download_image(image_prompt: str) -> str:
        """Search for an image and upload to S3.

        Args:
            image_prompt: Description of the image to search for

        Returns:
            S3 URI of the uploaded image with author attribution, or error message if failed.
            Format: "s3://bucket/key|author_name" on success, "error:reason" on failure.
        """
        if not unsplash_key:
            return "error:no_image_source"

        # 1. Search Unsplash
        response = requests.get(
            "https://api.unsplash.com/search/photos",
            params={"query": image_prompt, "orientation": "landscape", "per_page": 1},
            headers={"Authorization": f"Client-ID {unsplash_key}"},
        )

        if response.status_code != 200:
            return "error:search_failed"

        results = response.json().get("results", [])
        if not results:
            return "error:no_results"

        photo = results[0]
        image_url = photo["urls"]["regular"]
        author = photo["user"]["name"]

        # 2. Download image
        img_response = requests.get(image_url)
        if img_response.status_code != 200:
            return "error:download_failed"

        # 3. Upload to S3
        image_name = f"img_{generate(NANOID_ALPHABET, 8)}.jpg"
        s3_key = f"{artifact_base_path}/{image_name}"

        s3 = boto3.client("s3")
        s3.put_object(
            Bucket=bucket_name,
            Key=s3_key,
            Body=img_response.content,
            ContentType="image/jpeg",
        )

        # Return S3 URI with attribution
        return f"s3://{bucket_name}/{s3_key}|{author}"

    return search_and_download_image


@contextmanager
def get_report_agent(
    session_id: str,
    project_id: str | None = None,
    user_id: str | None = None,
):
    """Get a report agent instance with S3-based session management.

    The report agent specializes in creating PowerPoint presentations
    based on research and gathered information.

    Args:
        session_id: Unique identifier for the session
        project_id: Project ID (optional)
        user_id: User ID for session isolation (optional)

    Yields:
        Report agent instance with session management configured
    """

    config = get_config()
    interpreter = AgentCoreCodeInterpreter(
        region=config.aws_region,
        session_name=session_id,
        identifier=config.code_interpreter_identifier or None,
    )

    # Generate S3 base path for artifact
    artifact_id = f"art_{generate(NANOID_ALPHABET, 12)}"
    artifact_base_path = f"{user_id}/{project_id}/artifacts/{artifact_id}"
    bucket_name = config.agent_storage_bucket_name

    tools = [
        current_time,
        http_request,
        interpreter.code_interpreter,
        get_image_tool(bucket_name, artifact_base_path, config.unsplash_access_key),
    ]

    # Image handling guide
    image_guide = """
## Image Handling

### Before generating PPTX:
1. Parse the slide content to find `image_prompt` entries
2. For each image_prompt, call `search_and_download_image(prompt)`
3. Replace image_prompt with returned s3_uri in your slide data

### In code_interpreter:
Images are pre-downloaded to S3. Use the s3_uri directly:
```python
import boto3
from io import BytesIO

def load_image_from_s3(s3_uri: str) -> BytesIO:
    \"\"\"Load image from S3 URI.\"\"\"
    # Parse s3://bucket/key
    parts = s3_uri.replace("s3://", "").split("/", 1)
    bucket, key = parts[0], parts[1]

    s3 = boto3.client('s3')
    response = s3.get_object(Bucket=bucket, Key=key)
    return BytesIO(response['Body'].read())

# Usage
img_stream = load_image_from_s3(s3_uri)
slide.shapes.add_picture(img_stream, Inches(5.2), Inches(1.2), width=Inches(4.5))

# Add attribution (author name is after | in s3_uri result)
# e.g., "s3://bucket/key|Author Name" -> author = "Author Name"
```

### Fallback (if s3_uri starts with "error:"):
Create placeholder rectangle with gray background:
```python
from pptx.enum.shapes import MSO_SHAPE

# Create placeholder box instead of image
placeholder = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(x), Inches(y), Inches(width), Inches(height)
)
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
placeholder.line.fill.background()  # No border

# Add text describing intended image
tf = placeholder.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = f"[Image: {image_prompt}]"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(128, 128, 128)
p.alignment = PP_ALIGN.CENTER
```
"""

    system_prompt = f"""You are a Report Agent specialized in creating PowerPoint presentations.
{image_guide}

## S3 Upload Information
- **Bucket**: `{bucket_name}`
- **Base Path**: `{artifact_base_path}`
- **Filename**: Generate from presentation title (e.g., "AI Trends 2024" → "ai_trends_2024.pptx")

## CRITICAL: Generate ALL slides in a SINGLE code_interpreter call

You MUST generate the entire presentation in ONE code execution block.
DO NOT create slides one by one in separate tool calls.

## DO NOT:
- ❌ Create slides in multiple code_interpreter calls
- ❌ Ask user for approval between slides
- ❌ Generate slides incrementally

## DO:
- ✅ Plan all slides first (outline in your thinking)
- ✅ Generate complete PPTX in ONE code execution
- ✅ Upload once after completion

## Workflow

### Step 1: Research & Plan
- Use http_request to gather information if needed
- Create a mental outline of ALL slides before coding
- Plan: Title → Agenda → Content slides → Summary

### Step 2: Generate COMPLETE PPTX in ONE Script
Write a SINGLE Python script that creates ALL slides:

```python
!pip install python-pptx requests

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]

# ============ ALL SLIDES IN ONE SCRIPT ============

# --- Slide 1: Title ---
slide = prs.slides.add_slide(blank_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Presentation Title"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# --- Slide 2: Agenda ---
slide = prs.slides.add_slide(blank_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Agenda"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# --- Slides 3-N: Content (use loop for multiple similar slides) ---
topics = [
    {{"title": "Topic 1", "points": ["Point A", "Point B", "Point C"]}},
    {{"title": "Topic 2", "points": ["Point D", "Point E", "Point F"]}},
    {{"title": "Topic 3", "points": ["Point G", "Point H", "Point I"]}},
]

for topic in topics:
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = topic["title"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(topic["points"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {{point}}"
        p.font.size = Pt(20)
        p.space_before = Pt(12)

# --- Final Slide: Summary ---
slide = prs.slides.add_slide(blank_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Summary"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# ============ SAVE ============
prs.save('./presentation.pptx')
print(f"Created {{len(prs.slides)}} slides successfully")
```

### Step 3: Upload to S3 (in the SAME script)
After saving the PPTX file, upload directly to S3 using boto3:

```python
import boto3
import re

# Generate filename from title (lowercase, replace spaces with underscores, remove special chars)
def title_to_filename(title: str) -> str:
    filename = title.lower()
    filename = re.sub(r'[^a-z0-9\\s]', '', filename)
    filename = re.sub(r'\\s+', '_', filename.strip())
    return f"{{filename}}.pptx"

# Upload to S3
s3 = boto3.client('s3')
bucket = "{bucket_name}"
base_path = "{artifact_base_path}"
filename = title_to_filename(presentation_title)  # Use your presentation title
key = f"{{base_path}}/{{filename}}"

with open('./presentation.pptx', 'rb') as f:
    s3.upload_fileobj(
        f,
        bucket,
        key,
        ExtraArgs={{'ContentType': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'}}
    )

print(f"Uploaded to s3://{{bucket}}/{{key}}")
```

### Step 4: Report Success
After upload, report the S3 key to the user:
- S3 Key: `{artifact_base_path}/<generated_filename>.pptx`

## python-pptx Reference

**Positioning & Sizing:**
- `Inches(n)` - convert inches to EMUs
- `Pt(n)` - convert points to EMUs
- Slide: 10" x 5.625" (16:9)

**Colors:**
- `RGBColor(r, g, b)` - RGB values 0-255
- Common: Navy `(0,51,102)`, White `(255,255,255)`, Black `(0,0,0)`

**Text Alignment:**
- `PP_ALIGN.LEFT`, `PP_ALIGN.CENTER`, `PP_ALIGN.RIGHT`

**Font Properties:**
- `p.font.size = Pt(20)`
- `p.font.bold = True`
- `p.font.italic = True`
- `p.font.color.rgb = RGBColor(...)`
- `p.font.name = "Arial"`

**Shapes:**
```python
# Rectangle
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(...)

# Line
shape = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, x1, y1, x2, y2)
```

**Images:**
```python
slide.shapes.add_picture('image.png', Inches(1), Inches(1), width=Inches(3))
```

## Slidev Layout Implementation Guide

The write_agent outputs Slidev-style markdown. You MUST implement each layout type correctly:

### 1. title_slide - Opening slide
```python
slide = prs.slides.add_slide(blank_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Use theme primaryColor

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Main Title"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
tf = subtitle_box.text_frame
p = tf.paragraphs[0]
p.text = "Subtitle"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER
```

### 2. default - Standard content slide
```python
slide = prs.slides.add_slide(blank_layout)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Slide Title"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4))
tf = content_box.text_frame
tf.word_wrap = True
points = ["Point 1", "Point 2", "Point 3"]
for i, point in enumerate(points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"• {{point}}"
    p.font.size = Pt(20)
    p.space_before = Pt(12)
```

### 3. two_column - Split layout with left/right content
Markdown uses `<!-- left -->` and `<!-- right -->` to separate columns.
```python
slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Two Column Title"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# Left column (x: 0.5, width: 4.3)
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(4))
tf = left_box.text_frame
tf.word_wrap = True
left_points = ["Left point 1", "Left point 2"]
for i, point in enumerate(left_points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"• {{point}}"
    p.font.size = Pt(18)
    p.space_before = Pt(10)

# Right column (x: 5.2, width: 4.3)
right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(4))
tf = right_box.text_frame
tf.word_wrap = True
right_points = ["Right point 1", "Right point 2"]
for i, point in enumerate(right_points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"• {{point}}"
    p.font.size = Pt(18)
    p.space_before = Pt(10)
```

### 4. image_right - Content left, image right
```python
from pptx.enum.shapes import MSO_SHAPE
import boto3
from io import BytesIO

slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Image Right Title"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# Content on left (x: 0.5, width: 4.5)
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(4))
tf = content_box.text_frame
tf.word_wrap = True
# ... add content

# Helper function for loading images from S3
def load_image_from_s3(s3_uri: str) -> BytesIO:
    \"\"\"Load image from S3 URI.\"\"\"
    parts = s3_uri.replace("s3://", "").split("/", 1)
    bucket, key = parts[0], parts[1]
    s3 = boto3.client('s3')
    response = s3.get_object(Bucket=bucket, Key=key)
    return BytesIO(response['Body'].read())

# Helper function for adding image or placeholder
def add_image_or_placeholder(slide, s3_uri_with_author, image_prompt, x, y, width, height):
    \"\"\"Add image from S3 or create placeholder if unavailable.

    Args:
        s3_uri_with_author: Result from search_and_download_image tool (format: "s3://bucket/key|author" or "error:reason")
        image_prompt: Original image prompt for placeholder text
    \"\"\"
    if s3_uri_with_author.startswith("error:"):
        # Fallback: placeholder with description
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
        placeholder.line.fill.background()
        tf = placeholder.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Image: {{image_prompt}}]"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = PP_ALIGN.CENTER
        return False

    # Parse S3 URI and author
    s3_uri, author = s3_uri_with_author.split("|", 1)
    img_stream = load_image_from_s3(s3_uri)
    slide.shapes.add_picture(img_stream, Inches(x), Inches(y), width=Inches(width))

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.3), Inches(9), Inches(0.3))
    p = attr_box.text_frame.paragraphs[0]
    p.text = f"Photo by {{author}} on Unsplash"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(128, 128, 128)
    return True

# Image on right (x: 5.2, width: 4.5)
# s3_uri_with_author is the result from search_and_download_image tool
add_image_or_placeholder(slide, s3_uri_with_author, image_prompt, 5.2, 1.2, 4.5, 3)
```

### 5. image_left - Image left, content right
```python
slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
# ...

# Image on left (x: 0.3, width: 4.5)
add_image_or_placeholder(slide, s3_uri_with_author, image_prompt, 0.3, 1.2, 4.5, 3)

# Content on right (x: 5.0, width: 4.5)
content_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.3), Inches(4.5), Inches(4))
# ...
```

### 6. image_center - Centered image with title/caption
```python
slide = prs.slides.add_slide(blank_layout)

# Title at top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Image Center Title"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# Centered image (x: 2.5, width: 5)
add_image_or_placeholder(slide, s3_uri_with_author, image_prompt, 2.5, 1.2, 5, 2.8)

# Caption below image
caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
tf = caption_box.text_frame
p = tf.paragraphs[0]
p.text = "Caption text"
p.font.size = Pt(16)
p.alignment = PP_ALIGN.CENTER
```

### 7. comparison - Table-based comparison
Markdown uses table syntax for comparison data.
```python
slide = prs.slides.add_slide(blank_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Before vs After"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

# Table (rows, cols, left, top, width, height)
table_shape = slide.shapes.add_table(4, 3, Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
table = table_shape.table

# Header row
headers = ["Metric", "Before", "After"]
for col, header in enumerate(headers):
    cell = table.cell(0, col)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
    p = cell.text_frame.paragraphs[0]
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.font.size = Pt(16)

# Data rows
data = [("Cost", "$125K", "$75K"), ("Time", "2 hours", "15 min"), ("Uptime", "99.5%", "99.99%")]
for row_idx, (metric, before, after) in enumerate(data, start=1):
    table.cell(row_idx, 0).text = metric
    table.cell(row_idx, 1).text = before
    table.cell(row_idx, 2).text = after
    for col in range(3):
        p = table.cell(row_idx, col).text_frame.paragraphs[0]
        p.font.size = Pt(14)
```

### 8. quote - Featured quote/message
```python
slide = prs.slides.add_slide(blank_layout)

# Light background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

# Large quote mark
quote_mark = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(1), Inches(1.2))
tf = quote_mark.text_frame
p = tf.paragraphs[0]
p.text = "\\u201C"  # Opening quote character
p.font.size = Pt(120)
p.font.color.rgb = RGBColor(200, 200, 200)

# Quote text
quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
tf = quote_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The quote text goes here."
p.font.size = Pt(28)
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

# Attribution
attr_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
tf = attr_box.text_frame
p = tf.paragraphs[0]
p.text = "— Author Name"
p.font.size = Pt(18)
p.alignment = PP_ALIGN.RIGHT
```

### 9. end - Closing/thank you slide
```python
slide = prs.slides.add_slide(blank_layout)

# Dark background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)

# Thank you text
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Contact info
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
tf = contact_box.text_frame
contact_lines = ["Questions?", "email@company.com"]
for i, line in enumerate(contact_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
```

## Design Principles

**Color Palette:**
- Primary: Navy blue `(0, 51, 102)` for titles/headers
- Accent: Use complementary colors for emphasis
- Background: White or light gray for content slides
- Text: High contrast (dark on light, light on dark)

**Typography:**
- Title slides: 40-48pt
- Slide titles: 28-36pt
- Body text: 18-24pt
- Use Arial or Calibri for readability

**Layout:**
- Maximum 6 bullet points per slide
- Keep text concise and scannable
- Leave adequate margins (0.5" minimum)
- Consistent positioning across slides

## Important Notes
- Always use `blank_layout = prs.slide_layouts[6]` for full control
- Save file as `./presentation.pptx` before uploading
- Test that generated PPTX opens correctly
- Parse the layout from Slidev markdown frontmatter (e.g., `layout: two_column`)
- Apply theme colors from frontmatter: `primaryColor`, `accentColor`
"""

    bedrock_model = BedrockModel(
        model_id=REPORT_MODEL_ID,
        region_name=config.aws_region,
        boto_client_config=Config(
            read_timeout=300,
            connect_timeout=10,
            retries={"max_attempts": 3},
        ),
    )

    agent = Agent(
        model=bedrock_model,
        system_prompt=system_prompt,
        tools=tools,
    )

    yield agent


def _run_pptx_sync(
    session_id: str,
    project_id: str | None,
    user_id: str | None,
    instructions: str,
) -> str:
    """Run pptx agent synchronously (for use with asyncio.to_thread)."""
    with get_report_agent(session_id, project_id, user_id) as agent:
        result = agent(instructions)
        return str(result)


def create_pptx_tool(session_id: str, project_id: str | None, user_id: str | None):
    """Create a pptx agent tool bound to session context."""

    @tool
    async def pptx_agent(instructions: str) -> str:
        """Create a PowerPoint presentation based on the given instructions.

        Use this tool to:
        - Generate PowerPoint presentations from confirmed plans
        - Create slides with proper formatting and design

        Args:
            instructions: The confirmed plan and context for creating the presentation

        Returns:
            Result of presentation creation including download URL
        """
        return await asyncio.to_thread(
            _run_pptx_sync, session_id, project_id, user_id, instructions
        )

    return pptx_agent
